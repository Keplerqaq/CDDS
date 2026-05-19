#!/usr/bin/env python3
"""生成开题答辩 PPT —— 第3组：小猫钓鱼 + 制造业统计分析"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ---- 颜色 ----
BG_DARK  = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT   = RGBColor(0x64, 0xB5, 0xF6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xAA, 0xAA, 0xAA)
YELLOW   = RGBColor(0xFF, 0xD5, 0x4F)
GREEN    = RGBColor(0x81, 0xC7, 0x84)
RED      = RGBColor(0xE8, 0x6A, 0x6A)
ORANGE   = RGBColor(0xFF, 0xAB, 0x40)
TABLE_H  = RGBColor(0x2A, 0x2A, 0x42)
QS_COLOR = RGBColor(0x42, 0xA5, 0xF5)   # 快排蓝
SS_COLOR = RGBColor(0xFF, 0xA7, 0x26)   # 选择排序橙


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, y=Inches(0), h=Inches(1.1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y, W, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x2A, 0x42)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    return shape


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_multiline(slide, left, top, width, height, lines, size=20,
                  color=WHITE, spacing=1.3, font_name="微软雅黑",
                  align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, b, s, c = line, False, size, color
        else:
            text = line[0]
            b = line[1] if len(line) > 1 else False
            s = line[2] if len(line) > 2 else size
            c = line[3] if len(line) > 3 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(s)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(6 * spacing)
    return tf


def add_rounded_box(slide, left, top, width, height, fill_color=TABLE_H):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, color=GRAY):
    """画一个右箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color, text="",
             size=14, color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_bottom_bar(slide, text="数据结构课程设计 · 第3组 · 开题答辩"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), H - Inches(0.4), W, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x2A, 0x42)
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                       Inches(2), Inches(2.2), Inches(9.333), Pt(3),
                       ).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = ACCENT
slide.shapes[-1].line.fill.background()

add_textbox(slide, Inches(2), Inches(1.0), Inches(9.333), Inches(1.2),
            "数据结构课程设计", size=48, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(2.5), Inches(9.333), Inches(1.0),
            "小猫钓鱼游戏 + 制造业增加值统计分析系统",
            size=30, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(2), Inches(3.6), Inches(9.333), Inches(2.0), [
    ("计科 241/242 班 · 第 3 组", False, 22, GRAY),
    ("指导教师：周文峰", False, 20, GRAY),
    ("2026 年 5 月 19 日", False, 18, GRAY),
], align=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ============================================================
# Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "目  录")
add_bottom_bar(slide)

toc_items = [
    ("01", "选题的意义", "两个题目各自对应什么数据结构，综合训练价值"),
    ("02", "问题分析", "游戏规则、数据规模、七项功能需求、核心难点"),
    ("03", "存储结构与算法接口", "链式队列/栈、顺序表+索引数组、两种排序对比"),
    ("04", "可行性分析", "技术准备、风险应对"),
    ("05", "人员分工", "分工表"),
    ("06", "预期结果", "交付物"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.4) + Inches(0.9) * i
    add_textbox(slide, Inches(1.2), y, Inches(0.8), Inches(0.7),
                num, size=32, bold=True, color=ACCENT)
    add_textbox(slide, Inches(2.2), y + Inches(0.05), Inches(3.0), Inches(0.45),
                title, size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(2.2), y + Inches(0.5), Inches(8.0), Inches(0.4),
                desc, size=16, color=GRAY)

# ============================================================
# Slide 3: 选题的意义
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "一、选题的意义")
add_bottom_bar(slide)

# 左：小猫钓鱼
add_rounded_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.0))
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5.0), Inches(0.6),
            "小猫钓鱼游戏", size=26, bold=True, color=YELLOW)
add_multiline(slide, Inches(1.0), Inches(2.4), Inches(5.0), Inches(3.5), [
    ("核心数据结构：链式队列 + 链式栈", True, 18, WHITE),
    ("", False, 8, WHITE),
    ("队列（FIFO）：手牌出队头、收牌入队尾", False, 16, GRAY),
    ("栈（LIFO）：桌面压栈放牌、出栈收牌", False, 16, GRAY),
    ("标记数组：O(1) 查重，桌面每种牌面至多一张", False, 16, GRAY),
    ("", False, 8, WHITE),
    ("价值：将课堂上的队列和栈从\"知道概念\"推进到", False, 16, GRAY),
    ("\"独立设计和串联运行\"，理解 ADT 的设计思想", False, 16, GRAY),
])

# 右：制造业
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0))
add_textbox(slide, Inches(7.2), Inches(1.7), Inches(5.0), Inches(0.6),
            "制造业统计分析系统", size=26, bold=True, color=GREEN)
add_multiline(slide, Inches(7.2), Inches(2.4), Inches(5.0), Inches(3.5), [
    ("核心结构：顺序表 + 索引数组", True, 18, WHITE),
    ("", False, 8, WHITE),
    ("96 国 × 21 年真实面板数据（1999-2019）", False, 16, GRAY),
    ("七项功能：导入→查询→增速→两种排名→分析→保存", False, 16, GRAY),
    ("", False, 8, WHITE),
    ("价值：排序算法综合应用 + 索引数组的间接寻址思维", False, 16, GRAY),
    ("用真实数据感受中国制造业在全球格局中的变迁", False, 16, GRAY),
])

# ============================================================
# Slide 4: 问题分析 — 小猫钓鱼（含流程图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "二、问题分析 · 小猫钓鱼游戏")
add_bottom_bar(slide)

# 左侧：流程文字
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(4.0), Inches(0.5),
            "游戏流程", size=22, bold=True, color=ACCENT)
add_multiline(slide, Inches(0.8), Inches(1.9), Inches(5.2), Inches(4.5), [
    ("① 洗牌：N 张牌（1-9）随机打乱，平分给甲乙", False, 16, GRAY),
    ("② 甲先出一张牌 → 放到桌面", False, 16, GRAY),
    ("③ 交替出牌，每轮判断：", False, 16, GRAY),
    ("   ✅ 桌面无此牌 → 压栈，标记置 1", False, 15, GREEN),
    ("   ⚠️ 桌面有此牌 → 自己牌入队尾", False, 15, RED),
    ("      → 栈顶连续出栈入队尾，到匹配牌为止", False, 15, RED),
    ("④ 一方手牌为空 → 对手获胜", False, 16, GRAY),
    ("⑤ 回合上限防止无限对局", False, 16, GRAY),
])

# 右侧：流程图
fx = Inches(6.0)
fy = Inches(1.4)
bw = Inches(2.5)
bh = Inches(0.8)
gap = Inches(0.15)

# 出牌
add_rect(slide, fx + Inches(1.2), fy, bw, bh, TABLE_H, "出一张牌", 16, WHITE, True)
# 判重（菱形用矩形代替）
add_rect(slide, fx + Inches(1.2), fy + bh + gap, bw, bh,
         ACCENT, "桌上有此牌？", 15, WHITE, True)

# 无匹配 → 压栈
add_arrow(slide, fx + Inches(3.4), fy + bh + gap, Inches(0.5), bh, GREEN)
add_rect(slide, fx + Inches(3.9), fy + bh + gap, Inches(1.8), bh,
         RGBColor(0x2E, 0x7D, 0x32), "压栈 + 标记", 14, WHITE)

# 有匹配 → 收牌
add_arrow(slide, fx + Inches(1.2), fy + bh * 2 + gap * 2, Inches(0.0), Inches(0.01), GRAY)
add_rect(slide, fx + Inches(1.2), fy + bh * 3 + gap * 2, bw, bh * 1.3,
         RGBColor(0xC6, 0x28, 0x28), "收牌入队尾\n(栈顶出到匹配)", 13, WHITE)

# 判胜负
add_rect(slide, fx + Inches(1.2), fy + bh * 4.5 + gap * 3, bw, bh,
         TABLE_H, "手牌为空？", 15, WHITE, True)
add_arrow(slide, fx + Inches(3.4), fy + bh * 4.5 + gap * 3, Inches(0.5), bh, YELLOW)
add_rect(slide, fx + Inches(3.9), fy + bh * 4.5 + gap * 3, Inches(1.8), bh,
         YELLOW, "对手获胜", 14, BG_DARK, True)

# 连接箭头
add_textbox(slide, fx + Inches(1.2), fy + bh * 2.2, Inches(2.5), Inches(0.3),
            "是 ↓", size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(slide, fx + Inches(1.2), fy + bh * 4.2, Inches(2.5), Inches(0.3),
            "否 ↓", size=13, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, fx + Inches(2.5), fy + bh + Inches(0.2), Inches(1.0), Inches(0.3),
            "否 →", size=13, bold=True, color=GREEN, align=PP_ALIGN.LEFT)

# ============================================================
# Slide 5: 问题分析 — 制造业统计（含数据流图）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "二、问题分析 · 制造业增加值统计分析系统")
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.0), Inches(0.4),
            "数据规模：96 个国家 × 21 年（1999–2019）× 4 个收入等级", size=18, color=GRAY)

# 七项功能
funcs = [
    ("① 数据导入", "从 txt 批量读取 96 国数据到顺序表", "fscanf + strcmp 映射"),
    ("② 查询", "国家名 + 年份 → 增加值 + 增速", "遍历顺序表、全字匹配"),
    ("③ 增速计算", "growth_rate = (当年−上年)/上年", "上年为 0 则增速 = 0"),
    ("④ 增加值排名", "逐年全量按增加值降序排", "快速排序 O(n log n)"),
    ("⑤ 增速排名", "按收入等级分四组 → 组内降序排", "选择排序 O(n²)"),
    ("⑥ 增加值分析", "min / max / 均值 / 方差", "样本方差 n−1"),
    ("⑦ 保存", "排名结果写到两个 txt", "文件名派生自原始文件名"),
]

for i, (name, desc, note) in enumerate(funcs):
    y = Inches(1.8) + Inches(0.52) * i
    add_textbox(slide, Inches(0.8), y, Inches(2.2), Inches(0.35),
                name, size=14, bold=True, color=WHITE)
    add_textbox(slide, Inches(3.1), y, Inches(5.3), Inches(0.35),
                desc, size=13, color=GRAY)
    add_textbox(slide, Inches(8.6), y, Inches(4.0), Inches(0.35),
                note, size=12, color=GRAY)

# 底部数据流
flow_y = Inches(5.6)
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5.0), Inches(0.35),
            "数据流关系", size=18, bold=True, color=ACCENT)

flow_steps = [
    ("① 读", GREEN),
    ("② 查", ACCENT),
    ("③ 算", ACCENT),
    ("④⑤ 排", QS_COLOR),
    ("⑥ 析", YELLOW),
    ("⑦ 存", ORANGE),
]

fw = Inches(1.5)
fh = Inches(0.65)
for i, (label, fc) in enumerate(flow_steps):
    x = Inches(1.2) + Inches(2.0) * i
    add_rect(slide, x, flow_y, fw, fh, fc, label, 15, WHITE if fc != YELLOW else BG_DARK, True)
    if i < len(flow_steps) - 1:
        add_arrow(slide, x + fw + Inches(0.05), flow_y + Inches(0.1),
                  Inches(0.35), Inches(0.45), fc)

# ============================================================
# Slide 6: 索引数组原理（重点！可视化改进）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "二、核心难点 · 索引数组原理")
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.0), Inches(0.4),
            "教材约束：排序不能改变原始记录的物理顺序，也不能新建一个排好序的记录序列",
            size=17, bold=False, color=RED)

# --- 可视化：上下两行 ---
# 顶部：数据行 r[]（锁定）
ry = Inches(2.2)
add_textbox(slide, Inches(0.8), ry - Inches(0.4), Inches(3.0), Inches(0.35),
            "原始数据 r[]（绝对不动）", size=16, bold=True, color=WHITE)

data_colors = [TABLE_H, TABLE_H, TABLE_H, TABLE_H, TABLE_H]
data_labels = ["中国 5000", "美国 8000", "印度 2000", "日本 4000", "德国 3000"]
dw = Inches(2.0)
dh = Inches(0.7)
for i, (label, dc) in enumerate(zip(data_labels, data_colors)):
    x = Inches(0.8) + Inches(2.1) * i
    add_rect(slide, x, ry, dw, dh, dc, label, 16, WHITE)
    add_rect(slide, x, ry + dh, dw, Inches(0.35),
             RGBColor(0x15, 0x2A, 0x42), f"r[{i}]", 11, GRAY)

# 底部：索引行 idx[]
iy = Inches(4.0)
add_textbox(slide, Inches(0.8), iy - Inches(0.4), Inches(4.0), Inches(0.35),
            "索引数组 idx[]（排序前 → 排序后）", size=16, bold=True, color=YELLOW)

idx_labels_before = ["idx[0]=0", "idx[1]=1", "idx[2]=2", "idx[3]=3", "idx[4]=4"]
idx_labels_after  = ["idx[0]=1", "idx[1]=0", "idx[2]=3", "idx[3]=4", "idx[4]=2"]
idx_colors_after  = [GREEN, GREEN, GREEN, GREEN, GREEN]
idx_colors_before = [GRAY, GRAY, GRAY, GRAY, GRAY]

for i in range(5):
    x = Inches(0.8) + Inches(2.1) * i
    # 排序前
    add_rect(slide, x, iy, dw, dh * 0.5, TABLE_H, idx_labels_before[i], 12, GRAY)
    # 排序后
    add_rect(slide, x, iy + dh * 0.55, dw, dh * 0.5, GREEN, idx_labels_after[i], 14, WHITE, True)

# 映射箭头说明
add_rect(slide, Inches(0.8), Inches(5.3), Inches(11.3), Inches(1.4), TABLE_H)
add_multiline(slide, Inches(1.2), Inches(5.4), Inches(10.5), Inches(1.2), [
    ("排序后如何读取排名：", True, 17, WHITE),
    ("r[idx[0]] = r[1] = 美国 8000 → 第 1 名    |    r[idx[1]] = r[0] = 中国 5000 → 第 2 名", False, 15, GRAY),
    ("r[idx[2]] = r[3] = 日本 4000 → 第 3 名    |    r[idx[3]] = r[4] = 德国 3000 → 第 4 名", False, 15, GRAY),
    ("只交换 idx[] 中的整数下标，r[] 的物理顺序从未改变 ✓", True, 15, ACCENT),
])

# ============================================================
# Slide 7: 存储结构 — 小猫钓鱼（字号加大）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "三、存储结构及接口设计 · 小猫钓鱼")
add_bottom_bar(slide)

# 链式队列
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "链式队列（手牌）", size=22, bold=True, color=YELLOW)
code_q = """typedef struct QNode {
    int card;
    struct QNode *next;
} QNode;

typedef struct {
    QNode *front;   // 队头
    QNode *rear;    // 队尾
} LinkQueue;

void queue_init(LinkQueue *q);
void enqueue(LinkQueue *q, int card);
int  dequeue(LinkQueue *q);"""

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_q
p.font.size = Pt(15)
p.font.color.rgb = GREEN
p.font.name = "Courier New"

# 链式栈
add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
            "链式栈（桌面牌堆）", size=22, bold=True, color=YELLOW)
code_s = """typedef struct SNode {
    int card;
    struct SNode *next;
} SNode;

typedef struct {
    SNode *top;     // 栈顶
} LinkStack;

void stack_init(LinkStack *s);
void push(LinkStack *s, int card);
int  pop(LinkStack *s);"""

txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_s
p.font.size = Pt(15)
p.font.color.rgb = GREEN
p.font.name = "Courier New"

# 底部
add_rounded_box(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.3),
                fill_color=TABLE_H)
add_textbox(slide, Inches(1.2), Inches(5.4), Inches(11.0), Inches(0.4),
            "选型理由", size=18, bold=True, color=ACCENT)
add_multiline(slide, Inches(1.2), Inches(5.8), Inches(11.0), Inches(0.8), [
    ("数组队列 → 假溢出；链式队列 → 动态分配，无容量上限", False, 15, GRAY),
    ("数组栈   → 需预设容量；链式栈   → 按需分配，桌面牌数动态增减", False, 15, GRAY),
])

# ============================================================
# Slide 8: 存储结构 — 制造业（字号加大）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "三、存储结构及接口设计 · 制造业统计")
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(6.0), Inches(0.5),
            "顺序表 + 索引数组", size=22, bold=True, color=YELLOW)

code_rec = """typedef struct {
    char  country[30];
    int   country_type;    // 0=低 1=中低等 2=中高等 3=高
    float value_added[21];
    float growth_rate[21];
    int   year[21];
    int   index_va[21];    // 增加值排名
    int   index_gr[21];    // 增速排名（组内）
} RecType;

typedef struct {
    RecType r[MAXSIZE];    // r[0] 为第一个国家
    int index_l[], index_ml[],
        index_mh[], index_h[];
    int length;
} SqList;"""

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_rec
p.font.size = Pt(14)
p.font.color.rgb = GREEN
p.font.name = "Courier New"

add_textbox(slide, Inches(7.8), Inches(1.3), Inches(5.0), Inches(0.5),
            "核心算法接口", size=22, bold=True, color=YELLOW)

apis = [
    ("① void MVA_SqList_Read(L, filename)", ""),
    ("② void MVA_SqList_Search(L)", ""),
    ("③ void MVA_SqList_Calculate(L)", ""),
    ("④ void MVA_SqList_Sort_Va(L)", "快速排序"),
    ("    int partition(L, a, low, high, year)", ""),
    ("    void quick_sort_idx(L, a, low, high, year)", ""),
    ("⑤ void MVA_SqList_Sort_Gr(L)", "选择排序、分组"),
    ("    void group_sort_select(L, g, sz, res, yr)", ""),
    ("⑥ void MVA_SqList_Analyze(L)", ""),
    ("⑦ void MVA_SqList_Save(L, src_name)", ""),
]

for i, (api, note) in enumerate(apis):
    y = Inches(1.9) + Inches(0.48) * i
    c = ACCENT if any(kw in api for kw in ["④", "⑤"]) else GRAY
    add_textbox(slide, Inches(7.8), y, Inches(4.8), Inches(0.30),
                api, size=13, bold=("void" in api), color=c)
    if note:
        add_textbox(slide, Inches(11.5), y, Inches(1.5), Inches(0.30),
                    note, size=11, bold=True, color=YELLOW)

# ============================================================
# Slide 9: 两种排序算法对比（色块区分）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "三、核心算法 · 两种排序方式对比")
add_bottom_bar(slide)

# 表头
headers = ["", "增加值排名 ④", "增速排名 ⑤"]
col_w = [Inches(2.5), Inches(4.8), Inches(4.8)]
col_x = [Inches(0.8)]
for i in range(1, 3):
    col_x.append(col_x[-1] + col_w[i-1] + Inches(0.1))

# 表头底色
hdr_colors = [RGBColor(0x15, 0x2A, 0x42), RGBColor(0x15, 0x3A, 0x62), RGBColor(0x52, 0x2A, 0x12)]
for i, (hdr, w, x, hc) in enumerate(zip(headers, col_w, col_x, hdr_colors)):
    add_rounded_box(slide, x, Inches(1.5), w, Inches(0.55), fill_color=hc)
    add_textbox(slide, x, Inches(1.52), w, Inches(0.5),
                hdr, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("排序算法", "快速排序", "选择排序"),
    ("时间复杂度", "O(n log n) 平均 / O(n²) 最坏", "O(n²)"),
    ("空间复杂度", "O(log n) 递归栈", "O(1) 原地"),
    ("排序范围", "全局 96 国，逐年全量排", "分四组，各组内逐年排"),
    ("是否分组", "否（全量排名）", "是（先按收入等级分四组）"),
    ("稳定性", "不稳定", "不稳定"),
    ("选型理由", "数据量适中，平均效率最优", "分组后每组仅 10-50 国，O(n²) 差别不大"),
    ("结果存储", "index_va[year]", "index_gr[year]（组内排名）"),
]

QS_ROW = RGBColor(0x1A, 0x3A, 0x5A)   # 快排列蓝底
SS_ROW = RGBColor(0x5A, 0x3A, 0x1A)   # 选择排序列橙底

for r, row in enumerate(rows):
    y = Inches(2.15) + Inches(0.58) * r
    bg = TABLE_H if r % 2 == 0 else BG_DARK
    for c in range(3):
        # 数据列加色块
        col_bg = bg
        if c == 1:
            col_bg = TABLE_H if r % 2 == 0 else RGBColor(0x22, 0x32, 0x4A)
        elif c == 2:
            col_bg = TABLE_H if r % 2 == 0 else RGBColor(0x3A, 0x2A, 0x22)
        add_rounded_box(slide, col_x[c], y, col_w[c], Inches(0.52), fill_color=col_bg)
        add_textbox(slide, col_x[c] + Inches(0.15), y + Inches(0.02),
                    col_w[c] - Inches(0.2), Inches(0.48),
                    row[c], size=14, bold=(c == 0),
                    color=WHITE if c == 0 else GRAY,
                    align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)

# 底部说明
add_textbox(slide, Inches(0.8), Inches(6.9), Inches(12.0), Inches(0.4),
            "教材要求 ④ 和 ⑤ 必须使用不同排序算法 — 体现了对多种排序算法的掌握",
            size=16, bold=True, color=ACCENT)

# ============================================================
# Slide 10: 可行性分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "四、可行性分析")
add_bottom_bar(slide)

feasibility = [
    ("技术可行性", GREEN, [
        "已掌握全部所需数据结构：",
        "链表/栈/队列/顺序表（郝斌课程）",
        "排序：冒泡/插入/选择/快排/归并",
        "C语言：struct/typedef/文件I/O",
        "数据量可控：96国×21年=2016条",
        "O(n²)也毫秒级，无性能瓶颈",
        "参考代码可作验证基准",
    ]),
    ("风险与应对", RED, [
        "索引数组易出错",
        "→ 参考教材索引实现，逐步验证",
        "",
        "两种排序可能误用",
        "→ 独立函数封装，不同入口",
        "",
        "文件路径问题",
        "→ 数据文件和 main.c 同目录",
    ]),
]

for i, (title, accent, items) in enumerate(feasibility):
    x = Inches(1.5) + Inches(6.0) * i
    add_rounded_box(slide, x, Inches(1.4), Inches(5.0), Inches(5.2), fill_color=TABLE_H)
    add_textbox(slide, x + Inches(0.2), Inches(1.5), Inches(4.4), Inches(0.5),
                title, size=22, bold=True, color=accent)
    add_multiline(slide, x + Inches(0.2), Inches(2.1), Inches(4.4), Inches(4.0),
                  items, size=15, color=GRAY)

# ============================================================
# Slide 11: 人员分工与进度
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "五、人员分工")
add_bottom_bar(slide)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "人员分工", size=22, bold=True, color=ACCENT)

div_headers = ["姓名", "分工内容"]
div_rows = [
    ["Keplerqaq", "总体设计 + 题目一（队列/栈/博弈逻辑）\n+ 题目二（排序算法/索引数组/数据分析）\n+ 测试 + 实验报告"],
]
div_w = [Inches(2.0), Inches(9.5)]

for c, (hdr, w) in enumerate(zip(div_headers, div_w)):
    x = Inches(0.8) + sum(div_w[:c])
    add_rounded_box(slide, x, Inches(1.8), w, Inches(0.5),
                    fill_color=RGBColor(0x15, 0x2A, 0x42))
    add_textbox(slide, x, Inches(1.82), w, Inches(0.5),
                hdr, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for r, row in enumerate(div_rows):
    y = Inches(2.3) + Inches(0.55) * r
    for c, (cell, w) in enumerate(zip(row, div_w)):
        x = Inches(0.8) + sum(div_w[:c])
        bg = TABLE_H if r % 2 == 0 else BG_DARK
        add_rounded_box(slide, x, y, w, Inches(1.0), fill_color=bg)
        lines = cell.split("\n")
        for li, line in enumerate(lines):
            add_textbox(slide, x + Inches(0.15), y + Inches(0.05 + 0.3 * li),
                        w - Inches(0.2), Inches(0.3),
                        line, size=14, color=GRAY if li > 0 else WHITE,
                        bold=(li == 0))



# ============================================================
# Slide 12: 预期结果
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_bar(slide, "六、预期结果")
add_bottom_bar(slide)

results = [
    ("小猫钓鱼游戏",
     "终端交互程序，支持随机洗牌发牌、二人交替出牌、\n自动判定收牌和胜负、回合上限保护。"),
    ("制造业统计分析系统",
     "菜单驱动的终端程序，完整实现七项功能：\n导入 → 查询 → 增速计算 → 两种排名 → 分析 → 保存。"),
    ("输出文件",
     "制造业分析_inputdate_Sorted.txt（逐年全量增加值排名）\n"
     "制造业分析_inputdate_Grouped_Sorted.txt（逐年分组增速排名）"),
    ("实验报告",
     "含需求分析、概要设计、详细设计、\n测试用例与分析、源代码、运行截图、心得总结。"),
]

for i, (title, desc) in enumerate(results):
    y = Inches(1.5) + Inches(1.4) * i
    add_rounded_box(slide, Inches(1.0), y, Inches(11.3), Inches(1.2), fill_color=TABLE_H)
    add_textbox(slide, Inches(1.3), y + Inches(0.1), Inches(4.0), Inches(0.5),
                title, size=22, bold=True, color=ACCENT)
    add_multiline(slide, Inches(1.3), y + Inches(0.55), Inches(10.5), Inches(0.7),
                  desc.split("\n"), size=16, color=GRAY)

# ============================================================
# Slide 13: 致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                       Inches(2), Inches(2.8), Inches(9.333), Pt(3),
                       ).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = ACCENT
slide.shapes[-1].line.fill.background()

add_textbox(slide, Inches(2), Inches(1.5), Inches(9.333), Inches(1.2),
            "致  谢", size=52, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(2), Inches(3.2), Inches(9.333), Inches(2.5), [
    ("感谢周文峰老师的指导", False, 26, GRAY),
    ("", False, 10, GRAY),
    ("恳请各位老师批评指正", False, 22, GRAY),
], align=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ---- 保存 ----
output_path = "/Users/kepler/code/CDDS/开题答辩.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页")
