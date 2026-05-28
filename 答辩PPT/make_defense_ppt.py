#!/usr/bin/env python3
"""按答辩稿生成中期答辩 PPT（代码块无注释版）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x64, 0xB5, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xAA, 0xAA, 0xAA)
GREEN  = RGBColor(0x81, 0xC7, 0x84)
QS_C   = RGBColor(0x42, 0xA5, 0xF5)
SS_C   = RGBColor(0xFF, 0xA7, 0x26)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_bar(slide, text, y=Inches(0), h=Inches(1.1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, W, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x2A, 0x42)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_text(slide, text, left, top, width, height,
             size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def code_block(slide, code, left, top, width, height, size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(size)
    p.font.color.rgb = GREEN
    p.font.name = 'Courier New'
    return tf


# ===== Slide 1: 封面 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "数据结构课程设计", Inches(1), Inches(1.8), Inches(11), Inches(1),
         size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "小猫钓鱼游戏 + 制造业增加值统计分析系统",
         Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         size=30, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "第 3 组    |    李彦恒  陈召劲  云惟旺",
         Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         size=22, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "计科 242 班    |    指导教师：周文峰",
         Inches(1), Inches(4.6), Inches(11), Inches(0.6),
         size=22, color=GRAY, align=PP_ALIGN.CENTER)

# ===== Slide 2: 小猫钓鱼 — 数据结构 + 基础操作 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "一、小猫钓鱼 — 数据结构与基础操作")

code_block(s,
    "typedef struct QNode { int card; struct QNode *next; } QNode, *PQNode;\n"
    "typedef struct { PQNode front, rear; }   LinkQueue, *PLinkQueue;\n\n"
    "typedef struct SNode { int card; struct SNode *next; } SNode, *PSNode;\n"
    "typedef struct { PSNode top; }            LinkStack, *PLinkStack;",
    Inches(0.3), Inches(1.4), Inches(6.5), Inches(2.6), size=14)

code_block(s,
    "queue_init / enqueue / dequeue / queue_length / queue_print / queue_destroy\n"
    "stack_init  / push   / pop    / stack_print  / stack_destroy\n\n"
    "dequeue 和 pop 均返回 int 牌面值\n"
    "查重、收牌入队、匹配判断 都需要这个值\n\n"
    "shuffle:  for (i = n-1; i > 0; i--)  { j = rand() % (i+1);  swap; }\n"
    "deal_cards:  前 18 张给甲，后 18 张给乙",
    Inches(0.3), Inches(4.2), Inches(6.5), Inches(1.8), size=14)

add_text(s,
    "数据结构要点\n"
    "▸ 队列：队头出牌 + 队尾收牌\n"
    "▸ 栈：压栈放牌 + 出栈收牌\n"
    "▸ Stack 用包装结构体\n"
    "   消除二级指针调用\n\n"
    "操作函数要点\n"
    "▸ 出队/出栈返回 int\n"
    "   → flag查重需要值\n"
    "   → 收牌入队需要值\n"
    "   → 匹配判断需要值\n"
    "   取数据 + 释放节点一步完成\n\n"
    "▸ 洗牌：rand 随机位置交换\n"
    "▸ 发牌：一副牌平分给两人",
    Inches(7.2), Inches(1.4), Inches(5.8), Inches(5.5), size=17, color=WHITE)

# ===== Slide 3: 小猫钓鱼 — play_turn 出牌核心 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "二、小猫钓鱼 — 核心出牌函数 play_turn")

code_block(s,
    "int play_turn(PLinkQueue player, PLinkStack table, int flag[], char who[]) {\n"
    "    int card = dequeue(player);\n"
    "    printf(\"%s出%d，\", who, card);\n"
    "\n"
    "    if (flag[card - 1] == 1) {\n"
    "        printf(\"桌上已有 %d，收牌\", card);\n"
    "        enqueue(player, card);\n"
    "        while (1) {\n"
    "            int top = pop(table);\n"
    "            flag[top - 1] = 0;\n"
    "            enqueue(player, top);\n"
    "            if (top == card) break;\n"
    "        }\n"
    "    } else {\n"
    "        printf(\"无匹配，留在桌面\");\n"
    "        push(table, card);\n"
    "        flag[card - 1] = 1;\n"
    "    }\n"
    "    printf(\"\\n\");\n"
    "    return queue_is_empty(player);\n"
    "}",
    Inches(0.3), Inches(1.4), Inches(8.5), Inches(5.5), size=14)

add_text(s,
    "查重逻辑\n"
    "▸ flag[card-1]==1：桌上已有这张牌\n"
    "   ① 先把自己出的牌入队尾\n"
    "   ② while(1) 从栈顶收牌\n"
    "   ③ 每收一张清 flag\n"
    "   ④ top==card 时停止\n"
    "   → 两张相同牌之间的所有牌\n"
    "     都被收走\n\n"
    "▸ flag[card-1]!=1：无匹配\n"
    "   压栈 + flag 标 1\n\n"
    "游戏不变量\n"
    "▸ 桌面任意时刻每种牌面\n"
    "   最多一张\n\n"
    "返回值\n"
    "▸ queue_is_empty(player)\n"
    "   手牌空 → 对方获胜",
    Inches(9.2), Inches(1.4), Inches(3.8), Inches(5.5), size=16, color=WHITE)

# ===== Slide 4: 小猫钓鱼 — 主程序 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "三、小猫钓鱼 — 主程序 fish_game")

code_block(s,
    "void fish_game(void) {\n"
    "    srand((unsigned int)time(NULL));\n\n"
    "    int deck[36], deck_size = 36, idx = 0;\n"
    "    for (int v = 1; v <= 9; v++)\n"
    "        for (int i = 0; i < 4; i++)\n"
    "            deck[idx++] = v;\n"
    "    shuffle(deck, deck_size);\n\n"
    "    LinkQueue player_a, player_b;\n"
    "    queue_init(&player_a);  queue_init(&player_b);\n"
    "    deal_cards(deck, deck_size, &player_a, &player_b);\n\n"
    "    LinkStack table;  stack_init(&table);\n"
    "    int table_flag[9] = {0};\n\n"
    "    int turn = 0;\n"
    "    while (1) {\n"
    "        turn++;\n"
    "        printf(\"第%d轮：\", turn);\n"
    "        if (play_turn(&player_a, &table, table_flag, \"甲\")) {\n"
    "            printf(\"乙获胜！\\n\\n\"); break;\n"
    "        }\n"
    "        if (play_turn(&player_b, &table, table_flag, \"乙\")) {\n"
    "            printf(\"甲获胜！\\n\\n\"); break;\n"
    "        }\n"
    "        printf(\"\\n\");\n"
    "    }\n"
    "    stack_print(&table);  queue_print(获胜方);\n"
    "    queue_destroy(&player_a);  queue_destroy(&player_b);\n"
    "    stack_destroy(&table);\n"
    "}",
    Inches(0.3), Inches(1.4), Inches(8.0), Inches(5.8), size=13)

add_text(s,
    "主流程\n"
    "▸ srand(time) 保证随机\n"
    "▸ 生成 36 张牌\n"
    "▸ 洗牌 → 发牌\n"
    "▸ while(1) 交替出牌\n"
    "▸ play_turn 返空 → 对方赢\n"
    "▸ 销毁队列/栈释放内存\n\n"
    "细节\n"
    "▸ table_flag[9] 下标 0-8\n"
    "   flag[card-1] 减 1 对齐\n"
    "▸ 出牌返回值决定胜负\n"
    "   每次出牌都判空",
    Inches(8.6), Inches(1.4), Inches(4.4), Inches(5.5), size=18, color=WHITE)

# ===== Slide 5: 制造业 — 数据结构 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "四、制造业 — 数据结构定义")

code_block(s,
    "typedef struct {\n"
    "    char  country[30];        int   country_type;\n"
    "    float value_added[21];    float growth_rate[21];\n"
    "    int   index_va[21];       int   index_gr[21];\n"
    "} RecType;\n\n"
    "typedef struct {\n"
    "    RecType r[96];\n"
    "    int index_l[96],  index_ml[96];\n"
    "    int index_mh[96], index_h[96];\n"
    "    int count_l, count_ml, count_mh, count_h;\n"
    "    int length;          int growth_done;\n"
    "} SqList, *PSqList;",
    Inches(0.3), Inches(1.4), Inches(8.0), Inches(4.5), size=13)

add_text(s,
    "RecType 每条记录包含\n"
    "▸ 国家名 & 收入等级 (0-3)\n"
    "▸ 21 年增加值 & 增速\n"
    "▸ index_va：全球排名 (1-96)\n"
    "▸ index_gr：组内排名 (每组独立)\n\n"
    "SqList 容器包含\n"
    "▸ r[96] 顺序表本体\n"
    "▸ 四组分组下标 + 人数\n"
    "▸ length 实际记录数\n"
    "▸ growth_done 守卫\n"
    "   sort_gr & analyze 先查它\n"
    "   防增速未算就读垃圾值\n\n"
    "▸ = {0} 全零初始化\n"
    "   index_va[0]==0 哨兵可靠",
    Inches(8.6), Inches(1.4), Inches(4.4), Inches(5.8), size=16, color=WHITE)

# ===== Slide 6: 制造业 — 功能 1-3 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "五、制造业 — 数据导入 · 查询 · 增速计算")

add_text(s, "▶ 1. 数据导入 — MVA_SqList_Read",
         Inches(0.4), Inches(1.3), Inches(12), Inches(0.5),
         size=22, color=ACCENT, bold=True)
code_block(s,
    "while (n < MAXSIZE) {\n"
    "    if (fscanf(fp, \"%s %s\", L->r[n].country, temp_type) != 2) break;\n"
    "    if      (strcmp(temp_type, \"低收入国家\") == 0)   L->r[n].country_type = 0;\n"
    "    else if (strcmp(temp_type, \"中低等收入国家\") == 0) L->r[n].country_type = 1;\n"
    "    else if (strcmp(temp_type, \"中高等收入国家\") == 0) L->r[n].country_type = 2;\n"
    "    else if (strcmp(temp_type, \"高收入国家\") == 0)   L->r[n].country_type = 3;\n"
    "    for (int i = 0; i < YEARS; i++)\n"
    "        fscanf(fp, \"%f\", &L->r[n].value_added[i]);\n"
    "    n++;\n"
    "}",
    Inches(0.3), Inches(1.9), Inches(6.3), Inches(2.5), size=13)
add_text(s,
    "▸ fscanf \"%s %s\" 一次读国家名和等级\n"
    "▸ strcmp 中文标签映射为 0-3\n"
    "▸ for 循环读 21 个 float 进 value_added",
    Inches(0.4), Inches(4.6), Inches(6.2), Inches(1.8), size=14, color=GRAY)

add_text(s, "▶ 2. 查询 — MVA_SqList_Search",
         Inches(7.0), Inches(1.3), Inches(6), Inches(0.5),
         size=22, color=ACCENT, bold=True)
code_block(s,
    "for (i = 0; i < L->length; i++)\n"
    "    if (strcmp(name, L->r[i].country) == 0) break;\n"
    "if (i >= L->length) { printf(\"未找到该国。\\n\"); return; }\n"
    "int idx = year - 1999;\n"
    "printf(\"%s%d年: 增加值 = %.2f 亿美元\", name, year, L->r[i].value_added[idx]);\n"
    "if (idx > 0) printf(\", 增速 = %.2f%%\", L->r[i].growth_rate[idx] * 100);",
    Inches(7.0), Inches(1.9), Inches(6.0), Inches(2.3), size=13)
add_text(s, "▸ strcmp 遍历 + idx = year-1999",
    Inches(7.0), Inches(4.3), Inches(5), Inches(0.5), size=14, color=GRAY)

add_text(s,
    "▶ 3. 增速计算 — MVA_SqList_Calculate\n"
    "   growth_rate[0] = 0  (1999 年没有上年)\n"
    "   for (k = 1; k < 21; k++):  prev = value_added[k-1]\n"
    "       if (prev == 0)  growth_rate[k] = 0\n"
    "       else  growth_rate[k] = (value_added[k] - prev) / prev\n"
    "   growth_done = 1",
    Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.8), size=18, color=WHITE)

# ===== Slide 7: 制造业 — 功能 4-5 增加值排名 + 增速排名 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "六、制造业 — 增加值排名 (快排) · 增速排名 (选择排序)")

add_text(s, "增加值排名 · 索引数组 + 快速排序 O(n log n)",
         Inches(0.3), Inches(1.3), Inches(6.5), Inches(0.5),
         size=20, color=QS_C, bold=True)
code_block(s,
    "int idx_arr[96];\n"
    "for (i = 0; i < 96; i++) idx_arr[i] = i;\n"
    "quick_sort(L, idx_arr, 0, 95, year);\n"
    "for (rank = 0; rank < 96; rank++)\n"
    "    L->r[idx_arr[rank]].index_va[year] = rank + 1;\n\n"
    "int partition(PSqList L, int a[], int low, int high, int year) {\n"
    "    int pivot = a[low];\n"
    "    while (low < high) {\n"
    "        while (low < high && val(a[high]) <= val(pivot)) high--;\n"
    "        a[low] = a[high];\n"
    "        while (low < high && val(a[low])  >= val(pivot)) low++;\n"
    "        a[high] = a[low];\n"
    "    }\n"
    "    a[low] = pivot;  return low;\n"
    "}",
    Inches(0.3), Inches(1.9), Inches(6.8), Inches(3.2), size=11)
add_text(s,
    "▸ 不移动 r[] ，只交换 idx_arr 里的下标\n"
    "▸ 比较的是 r[idx_arr[i]].value_added\n"
    "▸ pivot = a[low] 存的是 r[] 下标值\n"
    "▸ 外层 while：别停太早\n"
    "▸ 内层 while：别跑太远(防越界)\n"
    "▸ verbose 控制是否打印排名表",
    Inches(0.3), Inches(5.3), Inches(6.8), Inches(1.8), size=14, color=GRAY)

add_text(s, "增速排名 · 分组 + 选择排序 O(n²)",
         Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.5),
         size=20, color=SS_C, bold=True)
code_block(s,
    "char *type_name[] = {\"低收入\",\"中低等\",\"中高等\",\"高收入\"};\n"
    "int  *groups[] = {L->index_l, L->index_ml, L->index_mh, L->index_h};\n"
    "int   sizes[]  = {L->count_l, L->count_ml, L->count_mh, L->count_h};\n\n"
    "void group_sort_select(PSqList L, int *group, int size,\n"
    "                       int *result, int year) {\n"
    "    for (int i = 0; i < size; i++) result[i] = group[i];\n"
    "    for (int i = 0; i < size-1; i++) {\n"
    "        int max = i;\n"
    "        for (int j = i+1; j < size; j++)\n"
    "            if (gr(result[j]) > gr(result[max])) max = j;\n"
    "        if (max != i) { swap(result[i], result[max]); }\n"
    "    }\n"
    "}",
    Inches(7.5), Inches(1.9), Inches(5.5), Inches(3.0), size=11)
add_text(s,
    "▸ switch(country_type) 分组\n"
    "▸ 三组并行数组统一驱动\n"
    "▸ 73 行复制粘贴 → 12 行循环\n"
    "▸ 教材要求两种不同排序算法",
    Inches(7.5), Inches(5.1), Inches(5.5), Inches(1.5), size=14, color=GRAY)

# ===== Slide 8: 制造业 — 功能 6-7 + 菜单 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "七、制造业 — 统计分析 · 导出 · 菜单主程序")

add_text(s, "▶ 6. 统计分析 — MVA_SqList_Analyze",
         Inches(0.3), Inches(1.3), Inches(6.5), Inches(0.5),
         size=20, color=ACCENT, bold=True)
add_text(s,
    "查找国家 → 一趟扫描同时找 min/max/sum\n"
    "avg = sum / 21\n"
    "二趟扫描求离差平方和 Σ(X-X̄)²\n"
    "var = sum_sq / 20   (贝塞尔修正 n-1)\n"
    "输出：增加值 & 增速的 min / max / 均值 / 方差",
    Inches(0.3), Inches(1.9), Inches(6.3), Inches(2.0), size=16, color=WHITE)

add_text(s, "▶ 7. 导出 — MVA_SqList_Save",
         Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.5),
         size=20, color=ACCENT, bold=True)
add_text(s,
    "排名不可能为 0 → = {0} 哨兵判断是否已排名\n"
    "拷贝文件名去 .txt → 拼接后缀\n\n"
    "核心：rank_idx 反向索引\n"
    "\"每国排第几\" 翻成 \"第几名是谁\"\n"
    "把排名当下标用\n"
    "fprintf 写入 _Sorted.txt\n"
    "       + _Grouped_Sorted.txt",
    Inches(7.5), Inches(1.9), Inches(5.5), Inches(2.0), size=16, color=WHITE)

add_text(s, "▶ 菜单 — MVA_Menu_Show          ▶ 主程序 — manufacturing_system",
         Inches(0.3), Inches(4.3), Inches(12.5), Inches(0.5),
         size=20, color=ACCENT, bold=True)
code_block(s,
    "int MVA_Menu_Show(void) {                        void manufacturing_system(void) {\n"
    "    printf(\"... 菜单 0-7 ...\");                    SqList L = {0};\n"
    "    int key;                                      char filename[100] = \"制造业分析_inputdate.txt\";\n"
    "    while (1) {                                   while (1) {\n"
    "        if (scanf(\"%d\",&key)==1 && key>=0 && key<=7)    int key = MVA_Menu_Show();\n"
    "            return key;                               switch (key) {\n"
    "        while(getchar()!='\\n');                          case 0: return;\n"
    "        printf(\"输入无效\");                             case 1: Read ...\n"
    "    }                                                  case 2: Search ...\n"
    "}                                                          ...\n"
    "                                                   }\n"
    "                                               }\n"
    "                                           }",
    Inches(0.3), Inches(4.9), Inches(12.5), Inches(2.4), size=12)

# ===== Slide 9: 致谢 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "感谢聆听", Inches(1), Inches(2.5), Inches(11), Inches(1),
         size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "恳请各位老师批评指正", Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         size=28, color=GRAY, align=PP_ALIGN.CENTER)

prs.save("答辩PPT_中期.pptx")
print("Done: 答辩PPT_中期.pptx")
