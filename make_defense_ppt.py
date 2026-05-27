#!/usr/bin/env python3
"""生成中期/终期答辩 PPT：8-10 分钟"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x64, 0xB5, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
YELLOW = RGBColor(0xFF, 0xD5, 0x4F)
GREEN = RGBColor(0x81, 0xC7, 0x84)
RED = RGBColor(0xE8, 0x6A, 0x6A)
ORANGE = RGBColor(0xFF, 0xAB, 0x40)
QS_COLOR = RGBColor(0x42, 0xA5, 0xF5)
SS_COLOR = RGBColor(0xFF, 0xA7, 0x26)

def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def title_bar(slide, text, y=Inches(0), h=Inches(1.1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, W, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x15, 0x2A, 0x42)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_text(slide, text, left, top, width, height, size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_code_block(slide, code, left, top, width, height, size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(size)
    p.font.color.rgb = GREEN
    p.font.name = 'Courier New'
    return tf

# ===== Slide 1: 封面 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "数据结构课程设计", Inches(1), Inches(1.8), Inches(11), Inches(1),
         size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "小猫钓鱼游戏 + 制造业增加值统计分析系统", Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         size=30, color=ACCENT, bold=False, align=PP_ALIGN.CENTER)
add_text(s, "第 3 组    |    李彦恒  陈召劲  云惟旺", Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         size=22, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "计科 242 班    |    指导教师：周文峰", Inches(1), Inches(4.6), Inches(11), Inches(0.6),
         size=22, color=GRAY, align=PP_ALIGN.CENTER)

# ===== Slide 2: 项目概述 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "一、项目概述")
add_text(s, "两个子系统，两套数据结构", Inches(0.5), Inches(1.4), Inches(12), Inches(0.6),
         size=28, color=ACCENT, bold=True)

add_text(s,
    "▎小猫钓鱼游戏\n"
    "   链式队列 (玩家手牌 FIFO)  +  链式栈 (桌面牌堆 LIFO)  +  flag 数组 O(1) 查重\n"
    "   核心逻辑：交替出牌 → 查重判断 → 收牌/压栈，游戏不变量：桌面每张牌面唯一\n\n"
    "▎制造业增加值统计分析系统\n"
    "   顺序表 (RecType[96])  +  索引数组排序  +  快速排序 & 选择排序\n"
    "   数据：96 国 × 21 年 (1999-2019)，按收入等级分 4 组\n"
    "   7 项功能：导入 → 查询 → 增速计算 → 增加值排名 → 增速排名 → 统计分析 → 保存",
    Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5), size=22, color=WHITE)

# ===== Slide 3: 小猫钓鱼 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "二、小猫钓鱼 — 链式队列 + 链式栈")
add_text(s, "核心数据结构 & 出牌/收牌逻辑", Inches(0.5), Inches(1.4), Inches(12), Inches(0.6),
         size=28, color=ACCENT, bold=True)

add_code_block(s,
    "typedef struct QNode { int card; struct QNode *next; } QNode, *PQNode;\n"
    "typedef struct { PQNode front, rear; } LinkQueue, *PLinkQueue;   // 手牌\n"
    "typedef struct { PSNode top; }          LinkStack, *PLinkStack;   // 桌面\n"
    "\n"
    "int play_turn(PLinkQueue player, PLinkStack table, int flag[], char who[]) {\n"
    "    int card = dequeue(player);                       // 出队头\n"
    "    if (flag[card] == 1) {                            // O(1) 查重\n"
    "        enqueue(player, card);                        // 收回自己出的牌\n"
    "        while (1) { int top = pop(table);             // 从栈顶收\n"
    "            flag[top] = 0; enqueue(player, top);      // 清标记，入队尾\n"
    "            if (top == card) break; }                 // 收到匹配牌\n"
    "    } else { push(table, card); flag[card] = 1; }     // 压栈，标记\n"
    "}",
    Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.5), size=16)

# ===== Slide 4: 制造业数据结构 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "三、制造业 — 数据结构全景")
add_text(s, "RecType (一条记录)  ×  SqList (顺序表容器)", Inches(0.5), Inches(1.4), Inches(12), Inches(0.6),
         size=28, color=ACCENT, bold=True)

add_code_block(s,
    "typedef struct {\n"
    "    char  country[30];   int country_type;         // 0低 1中低 2中高 3高\n"
    "    float value_added[21];  float growth_rate[21]; // 增加值 & 增速\n"
    "    int   index_va[21];     int   index_gr[21];     // 排名索引\n"
    "} RecType;\n"
    "\n"
    "typedef struct {\n"
    "    RecType r[96];                     // 顺序表本体，不动\n"
    "    int index_l[96], index_ml[96];     // 低收入/中低等 分组\n"
    "    int index_mh[96], index_h[96];     // 中高等/高收入 分组\n"
    "    int count_l, count_ml, count_mh, count_h;  // 各组人数\n"
    "    int length, growth_done;\n"
    "} SqList, *PSqList;",
    Inches(0.5), Inches(2.1), Inches(7), Inches(4.8), size=15)

add_text(s,
    "▸ 顺序表：O(1) 随机访问，支持索引数组排序\n"
    "▸ index_va[year]：存每个国家\"排第几名\"\n"
    "▸ index_l/ml/mh/h：按收入等级分组\n"
    "▸ growth_done：防未计算就排名",
    Inches(8), Inches(2.3), Inches(4.5), Inches(4), size=20, color=WHITE)

# ===== Slide 5: 索引数组排序 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "四、索引数组排序 — 核心原理")
add_text(s, "不移动原始数据，只交换下标", Inches(0.5), Inches(1.4), Inches(12), Inches(0.6),
         size=28, color=ACCENT, bold=True)

add_code_block(s,
    "int idx_arr[96];\n"
    "for (i = 0; i < 96; i++) idx_arr[i] = i;     // [0,1,2,...,95]\n"
    "quick_sort(L, idx_arr, 0, 95, year);          // 比较值，交换下标\n"
    "                                              // idx_arr[0]=第1名国家下标\n"
    "for (rank = 0; rank < 96; rank++)\n"
    "    L->r[idx_arr[rank]].index_va[year] = rank + 1;  // 回填名次",
    Inches(0.5), Inches(2.1), Inches(6), Inches(3), size=15)

# Visual explanation on right side
add_text(s,
    "▎正向 (idx_arr)：名次 → 国家\n"
    "   idx_arr[0] = 15  →  第1名 = r[15]\n\n"
    "▎反向 (rank_idx)：排名当下标\n"
    "   rank_idx[r[i].index_va - 1] = i\n"
    "   \"每国排第几\" 反转成 \"第几名是谁\"\n\n"
    "▎本质：\n"
    "   普通数组：下标 → 排名\n"
    "   反向索引：排名 → 下标\n\n"
    "▎两个索引：\n"
    "   index_va：全局排名 (1~96)\n"
    "   index_gr：组内排名 (每组独立从1起)",
    Inches(7.2), Inches(2.1), Inches(5.5), Inches(5), size=18, color=WHITE)

# ===== Slide 6: 两种排序算法 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "五、两种排序算法 — 快速排序 vs 选择排序")
add_text(s, "", Inches(0.5), Inches(1.2), Inches(12), Inches(0.1), size=10, color=WHITE)

# Quick Sort
add_text(s, "增加值排名：快速排序 O(n log n)", Inches(0.5), Inches(1.5), Inches(6), Inches(0.6),
         size=24, color=QS_COLOR, bold=True)
add_code_block(s,
    "int partition(PSqList L, int a[], int low, int high, int year) {\n"
    "    int pivot = a[low];\n"
    "    while (low < high) {\n"
    "        while (low < high && val(a[high]) <= val(pivot)) high--;\n"
    "        a[low] = a[high];                     // 右边不合格→左移\n"
    "        while (low < high && val(a[low]) >= val(pivot)) low++;\n"
    "        a[high] = a[low];                     // 左边不合格→右移\n"
    "    }\n"
    "    a[low] = pivot; return low;               // 基准归位\n"
    "}",
    Inches(0.5), Inches(2.1), Inches(6.2), Inches(3.8), size=15)

# Selection Sort
add_text(s, "增速排名：选择排序 O(n²) + 分组", Inches(7), Inches(1.5), Inches(6), Inches(0.6),
         size=24, color=SS_COLOR, bold=True)
add_code_block(s,
    "void group_sort_select(PSqList L, int *group, int size,\n"
    "                       int *result, int year) {\n"
    "    for (int i = 0; i < size-1; i++) {\n"
    "        int max = i;\n"
    "        for (int j = i+1; j < size; j++)\n"
    "            if (gr(result[j]) > gr(result[max])) max = j;\n"
    "        if (max != i) swap(&result[i], &result[max]);\n"
    "    }\n"
    "}",
    Inches(7), Inches(2.1), Inches(6), Inches(3), size=15)

add_text(s,
    "▸ 快排用于全局96国排名，效率优先\n"
    "▸ 选择排序每组20-30国，O(n²) 够用\n"
    "▸ 教材要求两种不同排序算法\n"
    "▸ verbose 参数：菜单打印，Save 静默",
    Inches(0.5), Inches(6.1), Inches(12), Inches(1), size=18, color=GRAY)

# ===== Slide 7: 功能演示 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "六、完整流程演示")
add_text(s, "Import → Search → Calculate → Rank → Analyze → Save", Inches(0.5), Inches(1.4), Inches(12), Inches(0.6),
         size=26, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(s,
    "▎ 1. 数据导入    fscanf 读 96 国 × 21 年，strcmp 映射收入等级\n"
    "▎ 2. 数据查询    输入国家名 + 年份，strcmp 遍历查找，输出增加值 & 增速\n"
    "▎ 3. 增速计算    (当年 - 上年) / 上年，除零保护，year=0 不计算\n"
    "▎ 4. 增加值排名  逐年快排 96 国，index_va 存名次\n"
    "▎ 5. 增速排名    按收入等级分 4 组，组内选择排序，index_gr 存组内名次\n"
    "▎ 6. 统计分析    某国 21 年的 min / max / 均值 / 方差 (贝塞尔修正 n-1)\n"
    "▎ 7. 结果保存    输出 _Sorted.txt + _Grouped_Sorted.txt",
    Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5), size=22, color=WHITE)

# ===== Slide 8: 踩坑与收获 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
title_bar(s, "七、踩坑与收获")
add_text(s, "", Inches(0.5), Inches(1.2), Inches(12), Inches(0.1), size=10, color=WHITE)

add_text(s,
    "▎ scanf 死循环\n"
    "   scanf(\"%d\") 失败时输入留在缓冲区 → 无限循环 \"输入无效\"\n"
    "   修复：scanf 返回值检查 + while(getchar()!='\\n') 清空残留\n\n"
    "▎ 中文对齐\n"
    "   printf 用字节算宽度 (UTF-8 中文 3B)，终端用显示列宽 (2列)\n"
    "   结论：混合中英文场景不用列对齐，改用 \"国家 — 数值\" 格式\n\n"
    "▎ 除零保护\n"
    "   增速公式 (当年-上年)/上年，上年可能为 0\n"
    "   修复：if (prev==0) growth_rate=0  else 正常计算\n\n"
    "▎ 未初始化数据\n"
    "   SqList 局部变量需 = {0}，否则 index_va[0] 为随机垃圾值\n"
    "   Save 的 \"自动排名\" 检测依赖此设计",
    Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.2), size=20, color=WHITE)

# ===== Slide 9: 致谢 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, "感谢聆听", Inches(1), Inches(2.5), Inches(11), Inches(1),
         size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "恳请各位老师批评指正", Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         size=28, color=GRAY, align=PP_ALIGN.CENTER)

prs.save("答辩PPT_中期.pptx")
print("Done: 答辩PPT_中期.pptx")
